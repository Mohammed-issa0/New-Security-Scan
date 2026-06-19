# -*- coding: utf-8 -*-
"""
Generate thesis Word document and PowerPoint presentation.
Usage: py thesis/generate.py
"""
from __future__ import annotations

import json
import re
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt

THESIS_DIR = Path(__file__).resolve().parent
SOURCES = THESIS_DIR / "sources"
DIAGRAMS = THESIS_DIR / "diagrams"
SCREENSHOTS = THESIS_DIR / "screenshots"
CONFIG_PATH = THESIS_DIR / "config.json"

CHAPTER_FILES = [
    "front-matter.md",
    "chapter-01.md",
    "chapter-02.md",
    "chapter-03.md",
    "chapter-04.md",
    "chapter-05.md",
    "chapter-06.md",
    "chapter-07.md",
    "chapter-08.md",
    "chapter-09.md",
    "references.md",
    "appendices.md",
]

AR_FONT = "Arial"
DIAGRAM_CAPTIONS = {
    "01-architecture.png": "شكل (4-1): الهندسة المعمارية العامة للمنصة",
    "02-auth-flow.png": "شكل (5-1): تدفق المصادقة (OTP + JWT)",
    "03-scan-lifecycle.png": "شكل (3-1): دورة حياة الفحص الأمني",
    "04-erd.png": "شكل (4-2): نموذج البيانات (ERD) — الكيانات الأساسية",
    "05-use-case.png": "شكل (3-2): أصحاب المصلحة والمنصة",
    "06-bff-proxy.png": "شكل (4-3): طبقة BFF Proxy",
    "07-tools.png": "شكل (2-1): أدوات الفحص السبع المدعومة",
    "08-sequence-scan.png": "شكل (4-4): مخطط تسلسل إنشاء فحص",
    "09-jira-oauth.png": "شكل (6-1): تكامل Jira OAuth",
    "10-plans.png": "شكل (6-2): حوكمة الباقات (Plan Governance)",
}


def load_config() -> dict:
    return json.loads(CONFIG_PATH.read_text(encoding="utf-8"))


def set_rtl_paragraph(paragraph):
    paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    p_pr = paragraph._p.get_or_add_pPr()
    bidi = OxmlElement("w:bidi")
    p_pr.append(bidi)


def set_rtl_run(run, bold=False, size=14, color=None):
    run.font.name = AR_FONT
    run.font.size = Pt(size)
    run.bold = bold
    if color:
        run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    rtl = OxmlElement("w:rtl")
    r_pr.append(rtl)
    lang = OxmlElement("w:lang")
    lang.set(qn("w:bidi"), "ar-SA")
    r_pr.append(lang)


def add_rtl_paragraph(doc, text, style=None, bold=False, size=14, align_center=False):
    p = doc.add_paragraph(style=style)
    if align_center:
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    else:
        set_rtl_paragraph(p)
    run = p.add_run(text)
    set_rtl_run(run, bold=bold, size=size)
    p.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
    p.paragraph_format.space_after = Pt(6)
    return p


def parse_markdown_table(lines):
    rows = []
    for line in lines:
        if not line.strip().startswith("|"):
            break
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        if all(set(c) <= set("-: ") for c in cells):
            continue
        rows.append(cells)
    return rows


def add_table(doc, rows):
    if not rows:
        return
    cols = max(len(r) for r in rows)
    table = doc.add_table(rows=len(rows), cols=cols)
    table.style = "Table Grid"
    for i, row in enumerate(rows):
        for j in range(cols):
            cell = table.rows[i].cells[j]
            text = row[j] if j < len(row) else ""
            cell.text = ""
            p = cell.paragraphs[0]
            set_rtl_paragraph(p)
            run = p.add_run(text)
            set_rtl_run(run, bold=(i == 0), size=11)
    doc.add_paragraph()


def add_image(doc, path: Path, caption: str | None = None, width_cm=15):
    if not path.exists():
        return
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run()
    run.add_picture(str(path), width=Cm(width_cm))
    if caption:
        add_rtl_paragraph(doc, caption, bold=True, size=12, align_center=True)


def render_markdown(doc, content: str, diagram_hooks: dict | None = None):
    diagram_hooks = diagram_hooks or {}
    lines = content.splitlines()
    i = 0
    in_code = False
    code_buf = []

    while i < len(lines):
        line = lines[i]

        if line.strip().startswith("```"):
            in_code = not in_code
            if not in_code and code_buf:
                p = doc.add_paragraph()
                set_rtl_paragraph(p)
                run = p.add_run("\n".join(code_buf))
                run.font.name = "Courier New"
                run.font.size = Pt(9)
                code_buf = []
            i += 1
            continue

        if in_code:
            code_buf.append(line)
            i += 1
            continue

        if line.startswith("# "):
            add_rtl_paragraph(doc, line[2:].strip(), bold=True, size=20)
            i += 1
            continue
        if line.startswith("## "):
            add_rtl_paragraph(doc, line[3:].strip(), bold=True, size=16)
            i += 1
            continue
        if line.startswith("### "):
            add_rtl_paragraph(doc, line[4:].strip(), bold=True, size=14)
            i += 1
            continue

        if line.strip().startswith("|"):
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i])
                i += 1
            add_table(doc, parse_markdown_table(table_lines))
            continue

        if line.strip().startswith("- "):
            bullets = []
            while i < len(lines) and lines[i].strip().startswith("- "):
                bullets.append(lines[i].strip()[2:])
                i += 1
            for b in bullets:
                add_rtl_paragraph(doc, f"• {b}", size=13)
            continue

        if line.strip().startswith("**") and line.strip().endswith("**"):
            add_rtl_paragraph(doc, line.strip().strip("*"), bold=True, size=13)
            i += 1
            continue

        hook_key = line.strip()
        if hook_key in diagram_hooks:
            add_image(doc, diagram_hooks[hook_key], DIAGRAM_CAPTIONS.get(diagram_hooks[hook_key].name))

        if line.strip():
            add_rtl_paragraph(doc, line.strip(), size=13)
        i += 1


def setup_document_styles(doc):
    section = doc.sections[0]
    section.page_height = Cm(29.7)
    section.page_width = Cm(21.0)
    section.top_margin = Cm(2.5)
    section.bottom_margin = Cm(2.5)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.5)


def add_cover_pages(doc, cfg):
    for _ in range(1):
        add_rtl_paragraph(doc, cfg["university"], bold=True, size=18, align_center=True)
        add_rtl_paragraph(doc, cfg["faculty"], size=14, align_center=True)
        add_rtl_paragraph(doc, cfg["department"], size=14, align_center=True)
        doc.add_paragraph()
        add_rtl_paragraph(doc, "أطروحة مشروع تخرج", bold=True, size=22, align_center=True)
        doc.add_paragraph()
        add_rtl_paragraph(doc, cfg["project_title_ar"], bold=True, size=20, align_center=True)
        add_rtl_paragraph(doc, cfg["project_title_en"], size=14, align_center=True)
        doc.add_paragraph()
        doc.add_paragraph()
        add_rtl_paragraph(doc, f"إعداد الطالب: {cfg['student_name']}", size=14, align_center=True)
        add_rtl_paragraph(doc, f"إشراف: {cfg['supervisor_name']}", size=14, align_center=True)
        add_rtl_paragraph(doc, f"العام الجامعي: {cfg['year']}", size=14, align_center=True)
        doc.add_page_break()

    add_rtl_paragraph(doc, "الإهداء", bold=True, size=18, align_center=True)
    doc.add_paragraph()
    add_rtl_paragraph(doc, cfg["dedication"], size=14)
    doc.add_page_break()

    add_rtl_paragraph(doc, "الشكر والتقدير", bold=True, size=18, align_center=True)
    doc.add_paragraph()
    add_rtl_paragraph(doc, cfg["acknowledgment"], size=14)
    doc.add_page_break()


def build_word(cfg):
    doc = Document()
    setup_document_styles(doc)
    add_cover_pages(doc, cfg)

    # Diagram placement hooks after specific chapter headings
    chapter_diagrams = {
        "chapter-02.md": [DIAGRAMS / "07-tools.png"],
        "chapter-03.md": [DIAGRAMS / "03-scan-lifecycle.png", DIAGRAMS / "05-use-case.png"],
        "chapter-04.md": [
            DIAGRAMS / "01-architecture.png",
            DIAGRAMS / "04-erd.png",
            DIAGRAMS / "06-bff-proxy.png",
            DIAGRAMS / "08-sequence-scan.png",
        ],
        "chapter-05.md": [DIAGRAMS / "02-auth-flow.png"],
        "chapter-06.md": [DIAGRAMS / "09-jira-oauth.png", DIAGRAMS / "10-plans.png"],
        "chapter-08.md": [],
    }

    for fname in CHAPTER_FILES:
        path = SOURCES / fname
        if not path.exists():
            continue
        content = path.read_text(encoding="utf-8")
        render_markdown(doc, content)
        for img in chapter_diagrams.get(fname, []):
            add_image(doc, img, DIAGRAM_CAPTIONS.get(img.name))
        if fname == "chapter-08.md":
            add_rtl_paragraph(doc, "لقطات شاشة من التطبيق", bold=True, size=16)
            if SCREENSHOTS.exists():
                for shot in sorted(SCREENSHOTS.glob("*.png")):
                    add_image(doc, shot, f"شكل: {shot.stem}", width_cm=14)
        if fname == "appendices.md":
            appendix_a = SOURCES / "appendix-a-endpoints.md"
            if appendix_a.exists():
                render_markdown(doc, appendix_a.read_text(encoding="utf-8"))
        doc.add_page_break()

    out = THESIS_DIR / "اطروحة_مشروع_التخرج.docx"
    doc.save(out)
    print(f"Word saved: {out}")
    return out


def add_slide_title(slide, title, subtitle=None):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PptRGB(15, 23, 42)
    title_box = slide.shapes.add_textbox(PptInches(0.5), PptInches(0.4), PptInches(9), PptInches(1.2))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = PptPt(32)
    p.font.bold = True
    p.font.color.rgb = PptRGB(34, 211, 238)
    p.alignment = PP_ALIGN.RIGHT
    if subtitle:
        sub = slide.shapes.add_textbox(PptInches(0.5), PptInches(1.3), PptInches(9), PptInches(0.6))
        stf = sub.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.size = PptPt(16)
        sp.font.color.rgb = PptRGB(200, 210, 220)
        sp.alignment = PP_ALIGN.RIGHT


def add_bullet_slide(prs, title, bullets, image: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, title)
    left = PptInches(0.6)
    top = PptInches(1.8)
    width = PptInches(8.8 if not image else 4.8)
    height = PptInches(5)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = PptPt(18)
        p.font.color.rgb = PptRGB(226, 232, 240)
        p.alignment = PP_ALIGN.RIGHT
    if image and image.exists():
        slide.shapes.add_picture(str(image), PptInches(5.5), PptInches(1.8), width=PptInches(4))


def build_pptx(cfg):
    prs = Presentation()
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(7.5)

    # 1 Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, cfg["project_title_ar"], f"{cfg['student_name']} — {cfg['year']}")

    slides_data = [
        ("محتويات العرض", [
            "المشكلة والأهداف",
            "الحل والمعمارية",
            "الأدوات والفحص",
            "الواجهة الأمامية",
            "الخلفية API",
            "الأمان والاختبار",
            "النتائج والمستقبل",
        ]),
        ("المشكلة", [
            "تشتت مخرجات أدوات الفحص الأمني",
            "غياب حوكمة الاستخدام والرصيد",
            "صعوبة الفحص المصدّق على التطبيقات المحمية",
            "تأخر تحويل الثغرات إلى إجراءات (Jira/PDF)",
        ]),
        ("الأهداف", [
            "منصة ويب موحّدة لفحوصات الأمن السيبراني",
            "دعم 7 أدوات صناعية مع إعدادات مرنة",
            "نظام باقات ورصيد ووقت تشغيل",
            "تكامل Jira وذكاء اصطناعي ولوحة إدارة",
        ]),
        ("الحل المقترح", [
            "Next.js Frontend + BFF Proxy",
            "REST API Backend (OpenAPI)",
            "عمال VPS لتنفيذ الفحوصات",
            "Webhooks لاستقبال النتائج",
        ], DIAGRAMS / "01-architecture.png"),
        ("التقنيات", [
            "Frontend: Next.js 14, React 18, TypeScript, TanStack Query",
            "Validation: Zod + react-hook-form",
            "i18n: next-intl (ar/en)",
            "Backend: REST API, JWT, Stripe, Jira OAuth",
            "Testing: Playwright E2E",
        ]),
        ("الهندسة المعمارية", [
            "Browser → Next.js → BFF → API → VPS",
            "طبقة endpoints.ts كعقد موحّد",
            "Proxy يمرّر JWT والترويسات",
        ], DIAGRAMS / "01-architecture.png"),
        ("أدوات الفحص", [
            "ZAP — فحص Web/API",
            "nmap — شبكات وخدمات",
            "ffuf — اكتشاف مسارات",
            "wpscan, sqlmap, xss, ssl",
            "مستويات عمق: light / deep / aggressive",
        ], DIAGRAMS / "07-tools.png"),
        ("دورة حياة الفحص", [
            "1. اختيار الباقة والهدف",
            "2. إعداد الفحص والمصادقة",
            "3. الطابور والتنفيذ على VPS",
            "4. النتائج والتقرير وJira",
        ], DIAGRAMS / "03-scan-lifecycle.png"),
        ("المصادقة", [
            "تسجيل/دخول → تحدي OTP",
            "تحقق OTP → JWT (access + refresh)",
            "تجديد تلقائي للجلسة",
            "حماية Admin عبر middleware",
        ], DIAGRAMS / "02-auth-flow.png"),
        ("نموذج الفحص", [
            "اختيار أداة وهدف وعمق",
            "إعداد Headers وToken وCookies",
            "إعدادات خاصة بكل أداة",
            "تقييد حسب الباقة تلقائياً",
        ], SCREENSHOTS / "07-scans-new-ar.png"),
        ("تفاصيل الفحص", [
            "تبويبات: نظرة عامة، أدوات، ثغرات، تقرير",
            "متابعة ETA لكل أداة",
            "إلغاء الفحص أثناء التنفيذ",
        ], SCREENSHOTS / "06-scans-list-ar.png"),
        ("الثغرات والتقارير", [
            "عرض الشدة وCVSS والتوصيات",
            "تصدير PDF",
            "توليد تقرير AI",
            "إنشاء تذاكر Jira",
        ]),
        ("نظام الباقات", [
            "رصيد فحوصات ووقت أقصى",
            "أدوات مسموحة وقيود عمق",
            "منع الفحص المصدّق أو brute-force حسب الباقة",
            "شراء رصيد إضافي عبر Stripe",
        ], DIAGRAMS / "10-plans.png"),
        ("تكامل Jira", [
            "OAuth مع Atlassian Cloud",
            "ربط مشاريع وأهداف",
            "إنشاء تذاكر من الثغرات",
            "إدارة مطوري Jira وأدوارهم",
        ], DIAGRAMS / "09-jira-oauth.png"),
        ("الذكاء الاصطناعي", [
            "اقتراح إعداد فحص (ai/scan-config)",
            "Guided Setup تفاعلي",
            "تقرير ما بعد الفحص",
        ], SCREENSHOTS / "01-landing-ar.png"),
        ("لوحة الإدارة", [
            "إدارة مستخدمين وفحوصات",
            "طابور VPS وإعادة ترتيب",
            "سجلات تدقيق وباقات",
        ]),
        ("الأمان", [
            "OTP + JWT + Refresh",
            "RBAC للمسارات الإدارية",
            "Plan Enforcement على Backend",
            "توقيع نطاق الفحص scopeSigned",
        ]),
        ("الاختبارات", [
            "Playwright E2E: login, target, scan, PDF",
            "Admin smoke: 7 مسارات",
            "Jira developer mapping (mocked)",
            "تحقق عقد OpenAPI ↔ endpoints.ts",
        ]),
        ("تحقيق الأهداف", [
            "✓ مصادقة OTP/JWT",
            "✓ 7 أدوات فحص",
            "✓ باقات وStripe",
            "✓ Jira + AI + Admin",
            "✓ عربي/إنجليزي",
        ]),
        ("نقاط القوة", [
            "أدوات صناعية حقيقية",
            "OpenAPI شامل (67+ endpoint)",
            "حوكمة باقات متقدمة",
            "تجربة مستخدم حديثة",
            "اختبارات E2E آلية",
        ]),
        ("التحديات", [
            "Backend منفصل عن مستودع Frontend",
            "اعتماد على خادم API خارجي",
            "لا دعم لحل CAPTCHA تلقائياً",
        ]),
        ("العمل المستقبلي", [
            "دمج Backend في Docker Compose",
            "CI/CD فحوص تلقائية",
            "Dashboard تحليلي للثغرات",
            "إشعارات للثغرات الحرجة",
        ]),
        ("الخاتمة", [
            "منصة متكاملة لتشغيل فحوصات الأمن",
            "جمع أدوات + حوكمة + تكاملات",
            "جاهزة للتوسع والتطوير",
        ]),
        ("الخلفية API", [
            "89 endpoint موثّقة في OpenAPI 3.0.1",
            "وحدات: Auth, Scans, Plans, Billing, Jira, AI, Admin",
            "Webhooks: Stripe + VPS",
            "امتداد متصفح: Extension Scan API",
        ]),
        ("لقطات من التطبيق", [
            "واجهة عربية/إنجليزية كاملة",
            "نموذج فحص متقدم مع قيود الباقة",
            "تكامل Jira وإعدادات OAuth",
        ], SCREENSHOTS / "01-landing-ar.png"),
        ("شكراً لحسن استماعكم", [
            "هل من أسئلة؟",
            f"{cfg['student_name']}",
            cfg["project_title_ar"],
        ]),
    ]

    for item in slides_data:
        title = item[0]
        bullets = item[1]
        image = item[2] if len(item) > 2 else None
        add_bullet_slide(prs, title, bullets, image)

    out = THESIS_DIR / "عرض_المناقشة.pptx"
    prs.save(out)
    print(f"PPT saved: {out}")
    return out


def main():
    cfg = load_config()
    build_word(cfg)
    build_pptx(cfg)


if __name__ == "__main__":
    main()
